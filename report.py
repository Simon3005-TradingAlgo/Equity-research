"""
report.py - professioneller PPTX-Bericht (Einzeltitel, ohne Anlageempfehlung)
=============================================================================
build_pptx(main, models, blended, consensus=None, peer_rs=None) -> bytes (PPTX)

Spiegelt den vollen Dashboard-Inhalt: Profil/Beschreibung, Entwicklung,
Abschluesse, Bewertung (DCF + Methoden), Dividende, Peer-Vergleich, Qualitaet.
Charts via matplotlib (kein Kaleido). Keine Kauf-/Verkaufsempfehlung.
"""
import io
from datetime import date
import numpy as np

NAVY = (31, 56, 100); TEAL = (46, 134, 171); GREEN = (30, 125, 50); RED = (183, 28, 28)
AMBER = (214, 137, 16); GREY = (90, 90, 90); LIGHT = (235, 240, 247)
MM = 1_000_000
C_NAVY, C_TEAL, C_GREEN, C_RED, C_AMBER = "#1F3864", "#2E86AB", "#1E7D32", "#B71C1C", "#D68910"

def _ok(x): return x is not None and not (isinstance(x, float) and np.isnan(x))
def _eur(x, ccy="", d=2): return "n/v" if not _ok(x) else f"{x:,.{d}f}{(' ' + ccy) if ccy else ''}"
def _pct(x, d=1): return "n/v" if not _ok(x) else f"{x*100:.{d}f}%"
def _xx(x): return "n/v" if not _ok(x) else f"{x:.1f}x"
def _arr(df, name): return df.loc[name].values.astype(float) if name in df.index else None

# ----------------------------------------------------------------- charts
def _mpl():
    import matplotlib; matplotlib.use("Agg"); import matplotlib.pyplot as plt; return plt

def _save(fig):
    buf = io.BytesIO(); fig.tight_layout(); fig.savefig(buf, format="png"); 
    _mpl().close(fig); buf.seek(0); return buf

def _fig_rev_margin(main, ccy):
    plt = _mpl(); yrs = [str(y) for y in main["years"]]
    rev = _arr(main["financials"], "Umsatz"); ebm = _arr(main["ratios"], "EBIT-Marge")
    fig, ax = plt.subplots(figsize=(6.0, 3.1), dpi=160)
    ax.bar(yrs, rev, color=C_NAVY, width=0.6); ax.set_ylabel(f"Umsatz (Mio. {ccy})", color=C_NAVY, fontsize=9)
    ax2 = ax.twinx(); ax2.plot(yrs, ebm * 100, color=C_TEAL, marker="o", lw=2); ax2.set_ylabel("EBIT-Marge %", color=C_TEAL, fontsize=9)
    ax.set_title("Umsatz & EBIT-Marge", fontsize=11, color=C_NAVY, loc="left")
    ax.spines["top"].set_visible(False); ax2.spines["top"].set_visible(False)
    ax.tick_params(labelsize=8); ax2.tick_params(labelsize=8); return _save(fig)

def _fig_eps_fcf(main, ccy):
    plt = _mpl(); yrs = [str(y) for y in main["years"]]
    eps = _arr(main["financials"], "EPS"); fcf = _arr(main["financials"], "Free Cashflow")
    fig, ax = plt.subplots(figsize=(6.0, 3.1), dpi=160)
    ax.bar(yrs, eps, color=C_NAVY, width=0.6); ax.set_ylabel(f"EPS ({ccy})", color=C_NAVY, fontsize=9)
    ax2 = ax.twinx(); ax2.plot(yrs, fcf, color=C_GREEN, marker="o", lw=2); ax2.set_ylabel(f"Free Cashflow (Mio. {ccy})", color=C_GREEN, fontsize=9)
    ax.set_title("EPS & Free Cashflow", fontsize=11, color=C_NAVY, loc="left")
    ax.spines["top"].set_visible(False); ax2.spines["top"].set_visible(False)
    ax.tick_params(labelsize=8); ax2.tick_params(labelsize=8); return _save(fig)

def _fig_leverage(main):
    plt = _mpl(); yrs = [str(y) for y in main["years"]]
    nde = _arr(main["ratios"], "Nettoverschuldung/EBITDA"); cov = _arr(main["ratios"], "Zinsdeckung")
    fig, ax = plt.subplots(figsize=(6.0, 3.1), dpi=160)
    ax.bar(yrs, nde, color=C_NAVY, width=0.6); ax.set_ylabel("Net Debt / EBITDA", color=C_NAVY, fontsize=9)
    ax2 = ax.twinx()
    if cov is not None: ax2.plot(yrs, cov, color=C_AMBER, marker="o", lw=2)
    ax2.set_ylabel("Zinsdeckung", color=C_AMBER, fontsize=9)
    ax.set_title("Verschuldung & Zinsdeckung", fontsize=11, color=C_NAVY, loc="left")
    ax.spines["top"].set_visible(False); ax2.spines["top"].set_visible(False)
    ax.tick_params(labelsize=8); ax2.tick_params(labelsize=8); return _save(fig)

def _fig_football(main, models, ccy):
    plt = _mpl(); price = main["latest"]["price"]
    items = [(k, v) for k, v in models.items() if _ok(v)]
    fig, ax = plt.subplots(figsize=(6.0, 3.1), dpi=160)
    cols = [C_GREEN if v >= (price or 0) else C_RED for _, v in items]
    ax.barh([k for k, _ in items], [v for _, v in items], color=cols)
    for i, (_, v) in enumerate(items): ax.text(v, i, f" {v:,.1f}", va="center", fontsize=8)
    if _ok(price):
        ax.axvline(price, color=C_NAVY, ls="--", lw=2)
        ax.text(price, len(items) - 0.4, f" Kurs {price:,.1f}", color=C_NAVY, fontsize=8)
    ax.set_title(f"Fairer Wert je Methode vs. Kurs ({ccy})", fontsize=11, color=C_NAVY, loc="left")
    ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
    ax.tick_params(labelsize=8); return _save(fig)

def _fig_dividend(main, ccy):
    plt = _mpl(); dh = main.get("dividends") or []
    by = {}
    for d in dh: by[d["year"]] = by.get(d["year"], 0) + d["amount"]
    yrs = sorted(by)[-4:]
    if not yrs: return None
    vals = [by[y] for y in yrs]
    fig, ax = plt.subplots(figsize=(6.0, 3.0), dpi=160)
    ax.bar([str(y) for y in yrs], vals, color=C_NAVY, width=0.6)
    for i, v in enumerate(vals): ax.text(i, v, f"{v:.2f}", ha="center", va="bottom", fontsize=8)
    ax.set_title(f"Dividende je Aktie ({ccy})", fontsize=11, color=C_NAVY, loc="left")
    ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
    ax.tick_params(labelsize=8); return _save(fig)

def _fig_peer_scatter(allr, names):
    plt = _mpl()
    xs = [r["headline"].get("roic") for r in allr]; ys = [r["multiples"].get("ev_ebit") for r in allr]
    fig, ax = plt.subplots(figsize=(6.0, 3.4), dpi=160)
    for i, (x, y, n) in enumerate(zip(xs, ys, names)):
        if not _ok(x) or not _ok(y): continue
        ax.scatter(x * 100, y, s=(150 if i == 0 else 70), color=(C_NAVY if i == 0 else C_TEAL), zorder=3)
        ax.annotate(n, (x * 100, y), fontsize=8, xytext=(4, 4), textcoords="offset points")
    ax.set_xlabel("ROIC % (Qualitaet)", fontsize=9); ax.set_ylabel("EV/EBIT (Bewertung)", fontsize=9)
    ax.set_title("Qualitaet vs. Bewertung", fontsize=11, color=C_NAVY, loc="left")
    ax.grid(alpha=0.25); ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
    ax.tick_params(labelsize=8); return _save(fig)

def _fig_peer_bars(allr, names):
    plt = _mpl()
    pe = [r["multiples"].get("pe") for r in allr]; ev = [r["multiples"].get("ev_ebit") for r in allr]
    x = np.arange(len(names)); w = 0.38
    fig, ax = plt.subplots(figsize=(6.0, 3.4), dpi=160)
    ax.bar(x - w / 2, [v if _ok(v) else 0 for v in pe], w, label="KGV", color=C_NAVY)
    ax.bar(x + w / 2, [v if _ok(v) else 0 for v in ev], w, label="EV/EBIT", color=C_TEAL)
    ax.set_xticks(x); ax.set_xticklabels(names, fontsize=8); ax.legend(fontsize=8)
    ax.set_title("Bewertung (niedriger = guenstiger)", fontsize=11, color=C_NAVY, loc="left")
    ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
    ax.tick_params(labelsize=8); return _save(fig)

# ----------------------------------------------------------------- builder
def build_pptx(main, models, blended, consensus=None, peer_rs=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    peer_rs = peer_rs or []
    ccy = main.get("price_ccy") or main.get("currency") or ""
    price = main["latest"]["price"]; h, m = main["headline"], main["multiples"]; d = main.get("dcf") or {}
    vals = [v for v in models.values() if _ok(v)]
    lo, hi = (min(vals), max(vals)) if vals else (np.nan, np.nan)
    mos = (blended / price - 1) if _ok(price) and _ok(blended) else np.nan

    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide(): return prs.slides.add_slide(blank)
    def tb(s, l, t, w, ht, text, size=14, bold=False, color=(0, 0, 0), align=PP_ALIGN.LEFT):
        b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(ht)); f = b.text_frame
        f.word_wrap = True; p = f.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = RGBColor(*color); r.font.name = "Arial"; return b
    def bar(s, color=NAVY): 
        sp = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.16))
        sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(*color); sp.line.fill.background()
    def head(s, title): bar(s); tb(s, 0.6, 0.32, 12.2, 0.7, title, 24, True, NAVY)
    def footer(s):
        tb(s, 0.6, 7.06, 12.2, 0.4, "Nur zu Informationszwecken - keine Anlageempfehlung. "
           "Datenquelle: Yahoo Finance / FMP (zu plausibilisieren).", 9, False, GREY)
    def table(s, l, t, w, ht, rows, header, fontsize=10, first_bold=False):
        tbl = s.shapes.add_table(len(rows) + 1, len(header), Inches(l), Inches(t), Inches(w), Inches(ht)).table
        for j, hd in enumerate(header):
            c = tbl.cell(0, j); c.text = str(hd); c.fill.solid(); c.fill.fore_color.rgb = RGBColor(*NAVY)
            p = c.text_frame.paragraphs[0]
            if p.runs: p.runs[0].font.color.rgb = RGBColor(255, 255, 255); p.runs[0].font.size = Pt(fontsize); p.runs[0].font.bold = True
        for i, rowv in enumerate(rows, 1):
            for j, v in enumerate(rowv):
                c = tbl.cell(i, j); c.text = str(v) if v not in (None, "") else " "
                p = c.text_frame.paragraphs[0]
                if p.runs:
                    p.runs[0].font.size = Pt(fontsize)
                    if first_bold and j == 0: p.runs[0].font.bold = True
        return tbl
    def df_rows(df, fmt):
        out = []
        for name in df.index:
            row = [name] + [fmt(v, name) for v in df.loc[name].values]
            out.append(row)
        return out

    # 1 - Titel & Bewertungsspanne
    s = slide(); bar(s)
    tb(s, 0.6, 1.1, 12, 1.0, main["name"], 32, True, NAVY)
    tb(s, 0.6, 2.15, 12, 0.5, f"{main['ticker']}  ·  {ccy}  ·  Equity Research  ·  {date.today().strftime('%d.%m.%Y')}", 15, False, GREY)
    tb(s, 0.6, 3.2, 7, 0.4, "Bewertungsspanne (alle Methoden)", 13, True, NAVY)
    tb(s, 0.6, 3.7, 7, 0.6, f"{_eur(lo, ccy)}   bis   {_eur(hi, ccy)}", 20, True, (0, 0, 0))
    tb(s, 0.6, 4.5, 7, 0.4, f"Median {_eur(blended, ccy)}   ·   Auf-/Abschlag vs. Kurs {_pct(mos, 0)}", 14, False, (0, 0, 0))
    tb(s, 8.3, 3.2, 4.3, 0.4, "Aktueller Kurs", 13, True, NAVY)
    tb(s, 8.3, 3.6, 4.3, 0.6, _eur(price, ccy), 22, True, (0, 0, 0))
    tb(s, 8.3, 4.6, 4.3, 0.4, f"Qualitaets-Score {main['conviction']:.1f} / 5,0", 14, False, GREY)
    footer(s)

    # 2 - Unternehmensprofil
    s = slide(); head(s, "Unternehmensprofil")
    p = main.get("profile") or {}
    facts = [["Sektor", p.get("sector") or "n/v"], ["Industrie", p.get("industry") or "n/v"],
             ["Sitz", f"{(p.get('city')+', ') if p.get('city') else ''}{p.get('country') or 'n/v'}"],
             ["Mitarbeiter", f"{int(p['employees']):,}" if p.get("employees") else "n/v"],
             ["Marktkap. (Mio.)", _eur(p["market_cap"] / MM, ccy, 0) if p.get("market_cap") else "n/v"],
             ["Beta", f"{p['beta']:.2f}" if p.get("beta") else "n/v"],
             ["52W-Spanne", f"{_eur(p.get('lo52'), '', 0)} – {_eur(p.get('hi52'), ccy, 0)}" if p.get("hi52") else "n/v"]]
    table(s, 0.6, 1.4, 5.2, 3.6, facts, ["Angabe", "Wert"], first_bold=True)
    desc = (p.get("summary") or "Keine Beschreibung verfuegbar.")[:1500]
    tb(s, 6.1, 1.4, 6.6, 5.3, desc, 11, False, (40, 40, 40))
    footer(s)

    # 3 - Entwicklung (Charts)
    s = slide(); head(s, "Entwicklung")
    s.shapes.add_picture(_fig_rev_margin(main, ccy), Inches(0.4), Inches(1.4), width=Inches(6.1))
    s.shapes.add_picture(_fig_eps_fcf(main, ccy), Inches(6.7), Inches(1.4), width=Inches(6.1))
    s.shapes.add_picture(_fig_leverage(main), Inches(0.4), Inches(4.2), width=Inches(6.1))
    kpis = [["ROIC", _pct(h["roic"])], ["ROIC - WACC", _pct(h["roic_wacc"])], ["EBIT-Marge", _pct(h["ebit_margin"])],
            ["FCF-Marge", _pct(h["fcf_margin"])], ["Umsatz-CAGR", _pct(h["rev_cagr"])],
            ["Cash Conversion", _pct(h["cash_conv"], 0)]]
    table(s, 7.0, 4.3, 5.7, 2.4, kpis, ["Kennzahl", "Wert"], 11, first_bold=True)
    footer(s)

    # 4 - Abschluesse
    stm = main.get("statements") or {}
    if stm:
        s = slide(); head(s, "Abschluesse (Mio., 4 Jahre)")
        yrs = [str(y) for y in main["years"]]
        def f_in(v, name):
            if not _ok(v): return "–"
            return f"{v:,.2f}" if (str(name).startswith("EPS") or str(name).startswith("Aktien")) else f"{v:,.0f}"
        def f_cf(v, name):
            if not _ok(v): return "–"
            return f"{v:,.2f}" if "je Aktie" in str(name) else f"{v:,.0f}"
        tb(s, 0.5, 1.25, 6, 0.3, "Gewinn- und Verlustrechnung", 12, True, NAVY)
        table(s, 0.5, 1.6, 6.4, 4.2, df_rows(stm["income"], f_in), ["Position"] + yrs, 9)
        tb(s, 7.1, 1.25, 6, 0.3, "Bilanz", 12, True, NAVY)
        table(s, 7.1, 1.6, 5.7, 3.2, df_rows(stm["balance"], f_in), ["Position"] + yrs, 9)
        tb(s, 7.1, 5.05, 6, 0.3, "Kapitalflussrechnung", 12, True, NAVY)
        table(s, 7.1, 5.4, 5.7, 1.3, df_rows(stm["cashflow"], f_cf), ["Position"] + yrs, 9)
        footer(s)

    # 5 - Bewertung
    s = slide(); head(s, "Bewertung")
    s.shapes.add_picture(_fig_football(main, models, ccy), Inches(0.4), Inches(1.4), width=Inches(6.3))
    dcf_rows = [["WACC", _pct(d.get("wacc"))], ["Wachstum (Ph.1)", _pct(d.get("growth"))],
                ["Terminales Wachstum", _pct(d.get("terminal"))], ["Basis-FCF (Mio.)", _eur((d.get("base_fcf") or np.nan) / MM, "", 0)],
                ["Enterprise Value (Mio.)", _eur((d.get("ev_dcf") or np.nan) / MM, "", 0)],
                ["Fair Value je Aktie", _eur(d.get("fair_value"), ccy)], ["Sicherheitsmarge", _pct(d.get("mos"), 0)],
                ["Reverse-DCF impl. g", _pct(main.get("reverse_growth"))]]
    table(s, 7.0, 1.4, 5.7, 3.0, dcf_rows, ["DCF (FCFF)", ""], 10, first_bold=True)
    mult_rows = [["KGV", _xx(m["pe"])], ["EV/EBIT", _xx(m["ev_ebit"])], ["EV/EBITDA", _xx(m["ev_ebitda"])],
                 ["KBV", _xx(m["pb"])], ["FCF-Rendite", _pct(m["fcf_yield"])], ["Div.-Rendite", _pct(m["div_yield"])]]
    table(s, 7.0, 4.7, 5.7, 2.0, mult_rows, ["Multiplikatoren", ""], 10, first_bold=True)
    meth = []
    for k, v in models.items():
        if not _ok(v): continue
        meth.append([k, _eur(v, ccy), _pct((v / price - 1) if _ok(price) else np.nan, 0)])
    meth.append(["Median", _eur(blended, ccy), _pct(mos, 0)])
    table(s, 0.4, 4.85, 6.3, 1.8, meth, ["Methode", "Fair Value", "Auf-/Abschlag"], 9)
    footer(s)

    # 6 - Dividende
    dfig = _fig_dividend(main, ccy); di = main.get("div_info") or {}
    if dfig or di:
        s = slide(); head(s, "Dividende")
        if dfig: s.shapes.add_picture(dfig, Inches(0.4), Inches(1.5), width=Inches(6.6))
        dy = di.get("yield_")
        if dy is not None and dy > 0.5: dy = dy / 100
        drows = [["Dividendenrendite", _pct(dy)], ["Ausschuettungsquote", _pct(di.get("payout"))],
                 ["Letzter Ex-Tag", di.get("ex_date") or "n/v"], ["DPS (letztes GJ)", _eur(main["latest"].get("dps"), ccy)]]
        table(s, 7.3, 1.6, 5.4, 1.9, drows, ["Kennzahl", "Wert"], 11, first_bold=True)
        dh = main.get("dividends") or []
        if dh:
            recent = [[x["date"], f"{x['amount']:.3f} {ccy}"] for x in dh[::-1][:8]]
            table(s, 7.3, 3.9, 5.4, 2.8, recent, ["Ex-Tag", "Betrag"], 10)
        footer(s)

    # 7 - Peer-Vergleich
    if peer_rs:
        s = slide(); head(s, "Peer-Vergleich")
        allr = [main] + peer_rs; names = [r["ticker"] for r in allr]
        s.shapes.add_picture(_fig_peer_bars(allr, names), Inches(0.4), Inches(1.4), width=Inches(6.1))
        s.shapes.add_picture(_fig_peer_scatter(allr, names), Inches(6.7), Inches(1.4), width=Inches(6.1))
        rows = []
        for i, r in enumerate(allr):
            hh, mm = r["headline"], r["multiples"]
            nm = ("> " + r["ticker"]) if i == 0 else r["ticker"]
            rows.append([nm, f"{r['conviction']:.1f}", _pct(hh["mos"], 0), _pct(hh["roic"]),
                         _pct(hh["ebit_margin"]), _xx(hh["nd_ebitda"]), _xx(mm["pe"]), _xx(mm["ev_ebit"])])
        table(s, 0.4, 4.7, 12.5, 2.0, rows,
              ["Titel", "Score", "Auf/Ab", "ROIC", "EBIT-M.", "ND/EBITDA", "KGV", "EV/EBIT"], 9, first_bold=True)
        footer(s)

    # 8 - Qualitaets-Scorecard
    s = slide(); head(s, "Qualitaets-Scorecard")
    labels = dict(geschaeftsmodell="Geschaeftsmodell*", management="Management*", wachstum="Wachstum",
                  profitabilitaet="Profitabilitaet (ROIC>WACC)", bilanz="Bilanz & Verschuldung",
                  cashflow="Cashflow-Qualitaet", margen="Margen", bewertung="Bewertung (Niveau)")
    rows = [[labels[k], _pct(main["weights"][k], 0), str(main["scores"][k])] for k in main["weights"]]
    rows.append(["Qualitaets-Score (gewichtet)", "", f"{main['conviction']:.1f} / 5"])
    table(s, 0.6, 1.4, 8.0, 4.6, rows, ["Kriterium", "Gewicht", "Score (1-5)"], 11, first_bold=True)
    tb(s, 8.9, 1.5, 3.9, 4.0,
       "Der Qualitaets-Score bewertet ausschliesslich fundamentale Merkmale (Profitabilitaet, Bilanz, "
       "Cashflow, Wachstum, Bewertungsniveau). Einordnung der Unternehmensqualitaet, keine Kauf-/"
       "Verkaufsempfehlung.\n\n* qualitativ gesetzt.", 11, False, GREY)
    footer(s)

    bio = io.BytesIO(); prs.save(bio); return bio.getvalue()

if __name__ == "__main__":
    import equity_engine as eng
    r = eng.compute(eng._synthetic())
    r["dividends"] = [{"date": f"{y}-08-15", "year": y, "amount": a} for y, a in [(2022, 1.0), (2023, 1.1), (2024, 1.25), (2025, 1.4)]]
    r["div_info"] = {"yield_": 0.72, "payout": 0.28, "ex_date": "2025-05-15"}
    peer = eng.compute(eng._synthetic()); peer["ticker"] = "PEER1"
    models = {"DCF (FCFF)": r["headline"]["fair_value"], "DDM (Gordon)": 22.9, "EV/EBIT (Peer)": 49.4, "Analysten-Ziel": 42.0}
    open("selftest_report.pptx", "wb").write(build_pptx(r, models, 33.4, dict(target_mean=42.0), [peer]))
    print("PPTX ok")
